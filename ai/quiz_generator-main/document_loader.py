import re
from nltk.tokenize import sent_tokenize
import PyPDF2
from pptx import Presentation

class DocumentLoader:
    """Load and process PDF/PPT documents"""
    
    def __init__(self):
        try:
            from nltk.corpus import stopwords
            self.stop_words = set(stopwords.words('english'))
        except:
            self.stop_words = {'the', 'a', 'an', 'and', 'or', 'but', 'in', 'on', 'at', 'to', 
                              'for', 'of', 'with', 'by', 'is', 'are', 'was', 'were', 'be', 'been'}
    
    def load_pdf(self, file_path: str) -> str:
        """Load text from PDF file"""
        text = ""
        
        try:
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for page in pdf_reader.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text += page_text + "\n\n"
        except Exception as e:
            print(f"Warning: Could not extract text from PDF: {e}")
        
        return text
    
    def load_ppt(self, file_path: str) -> str:
        """Load text from PPT file"""
        text = ""
        
        try:
            presentation = Presentation(file_path)
            for slide in presentation.slides:
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_text.append(shape.text.strip())
                if slide_text:
                    text += "\n".join(slide_text) + "\n\n"
        except Exception as e:
            print(f"Warning: Could not extract text from PPT: {e}")
        
        return text
    
    def process_text(self, text: str):
        """Process extracted text into sentences and paragraphs"""
        # Clean text
        text = re.sub(r'\s+', ' ', text).strip()
        
        # Split into sentences and paragraphs
        sentences = [sent.strip() for sent in sent_tokenize(text) if len(sent.split()) > 3]
        paragraphs = [p.strip() for p in text.split('\n\n') if p.strip()]
        
        return text, sentences, paragraphs